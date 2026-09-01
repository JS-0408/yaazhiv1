"""
ingest_docs.py — Yaazhi Bulk Document Ingestion Script
Loads PDFs, DOCX, TXT, and PPTX files from the knowledge/ vault,
chunks them, embeds with Sentence-Transformers, and stores in ChromaDB / pgvector.

Usage:
    python scripts/ingest_docs.py
    python scripts/ingest_docs.py --source knowledge/btech_notes
    python scripts/ingest_docs.py --source knowledge/ieee_papers --category research
    python scripts/ingest_docs.py --reindex        # wipe and rebuild entire index
"""

import argparse
import hashlib
import json
import os
import sys
import time
from datetime import datetime
from pathlib import Path
from typing import Optional

# ─────────────────────────────────────────────────────────
# Allow running from repo root without installing the package
# ─────────────────────────────────────────────────────────
sys.path.insert(0, str(Path(__file__).parent.parent))

# ─────────────────────────────────────────────────────────
# Lazy imports — only fail at runtime if a dep is missing
# ─────────────────────────────────────────────────────────
def _import_or_die(module: str, pip_name: str):
    try:
        return __import__(module)
    except ImportError:
        print(f"❌  Missing dependency: '{pip_name}'.  Run: pip install {pip_name}")
        sys.exit(1)


# ─────────────────────────────────────────────────────────
# Constants
# ─────────────────────────────────────────────────────────
REPO_ROOT      = Path(__file__).parent.parent
KNOWLEDGE_DIR  = REPO_ROOT / "knowledge"
INDEX_FILE     = KNOWLEDGE_DIR / "index.json"

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".md", ".pptx", ".csv"}

# Chunk settings
CHUNK_SIZE    = 512   # tokens / approx chars (adjust for your embedding model)
CHUNK_OVERLAP = 64    # overlap to preserve context across chunk boundaries

# Category auto-detection from folder names
FOLDER_CATEGORIES = {
    "btech_notes": "academic",
    "ieee_papers":  "research",
    "projects":     "project",
}


# ─────────────────────────────────────────────────────────
# Argument parser
# ─────────────────────────────────────────────────────────
def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Ingest documents from the knowledge vault into Yaazhi's memory.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument(
        "--source",
        default=str(KNOWLEDGE_DIR),
        help="Path to document folder (default: knowledge/)",
    )
    parser.add_argument(
        "--category",
        default=None,
        help="Override category tag for all ingested documents",
    )
    parser.add_argument(
        "--reindex",
        action="store_true",
        help="Wipe the existing vector collection and rebuild from scratch",
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Scan and print files to be ingested without actually indexing",
    )
    parser.add_argument(
        "--chunk-size",
        type=int,
        default=CHUNK_SIZE,
        help=f"Characters per chunk (default: {CHUNK_SIZE})",
    )
    parser.add_argument(
        "--verbose",
        action="store_true",
        help="Print chunk-level detail",
    )
    return parser.parse_args()


# ─────────────────────────────────────────────────────────
# File discovery
# ─────────────────────────────────────────────────────────
def discover_files(source: str) -> list[Path]:
    """Recursively find all supported documents under source."""
    source_path = Path(source)
    if not source_path.exists():
        print(f"❌  Source directory does not exist: {source_path}")
        sys.exit(1)

    files = [
        f for f in source_path.rglob("*")
        if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS
    ]
    return sorted(files)


def file_hash(path: Path) -> str:
    """SHA-256 hash of file contents — used to skip already-indexed files."""
    h = hashlib.sha256()
    with open(path, "rb") as fh:
        for chunk in iter(lambda: fh.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def infer_category(path: Path, override: Optional[str]) -> str:
    """Determine category from folder structure or explicit override."""
    if override:
        return override
    for part in path.parts:
        if part in FOLDER_CATEGORIES:
            return FOLDER_CATEGORIES[part]
    return "general"


# ─────────────────────────────────────────────────────────
# Text extraction
# ─────────────────────────────────────────────────────────
def extract_text(path: Path) -> str:
    """Extract raw text from a document file."""
    ext = path.suffix.lower()

    if ext in {".txt", ".md", ".csv"}:
        return path.read_text(encoding="utf-8", errors="replace")

    if ext == ".pdf":
        try:
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                return "\n\n".join(
                    page.extract_text() or "" for page in pdf.pages
                )
        except ImportError:
            # Fallback to PyMuPDF
            fitz = _import_or_die("fitz", "pymupdf")
            doc = fitz.open(str(path))
            return "\n\n".join(page.get_text() for page in doc)

    if ext == ".docx":
        docx = _import_or_die("docx", "python-docx")
        from docx import Document
        doc = Document(str(path))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if ext == ".pptx":
        pptx = _import_or_die("pptx", "python-pptx")
        from pptx import Presentation
        prs = Presentation(str(path))
        slides_text = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            slides_text.append("\n".join(slide_text))
        return "\n\n---\n\n".join(slides_text)

    return ""


# ─────────────────────────────────────────────────────────
# Chunking
# ─────────────────────────────────────────────────────────
def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Split text into overlapping chunks."""
    if not text.strip():
        return []

    chunks = []
    start  = 0
    length = len(text)

    while start < length:
        end = min(start + chunk_size, length)

        # Try to break at a paragraph or sentence boundary
        if end < length:
            for sep in ["\n\n", "\n", ". ", " "]:
                boundary = text.rfind(sep, start, end)
                if boundary > start:
                    end = boundary + len(sep)
                    break

        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)

        start = max(start + 1, end - overlap)

    return chunks


# ─────────────────────────────────────────────────────────
# Index management (index.json)
# ─────────────────────────────────────────────────────────
def load_index() -> dict:
    if INDEX_FILE.exists():
        with open(INDEX_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    return {"version": 1, "documents": {}}


def save_index(index: dict) -> None:
    KNOWLEDGE_DIR.mkdir(parents=True, exist_ok=True)
    index["last_updated"] = datetime.utcnow().isoformat() + "Z"
    with open(INDEX_FILE, "w", encoding="utf-8") as f:
        json.dump(index, f, indent=2, ensure_ascii=False)


# ─────────────────────────────────────────────────────────
# Vector store wrapper
# ─────────────────────────────────────────────────────────
def get_vector_store():
    """Return a configured VectorStore instance, or None if unavailable."""
    try:
        from memory.vector_store import VectorStore
        return VectorStore()
    except Exception as exc:
        print(f"⚠️   Vector store unavailable: {exc}")
        print("     Chunks will be indexed in index.json only (no semantic search).")
        return None


# ─────────────────────────────────────────────────────────
# Main ingestion logic
# ─────────────────────────────────────────────────────────
def ingest(args: argparse.Namespace) -> None:
    print("\n━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━")
    print("  📚  Yaazhi Document Ingestion")
    print(f"  Source  : {args.source}")
    print(f"  Mode    : {'DRY RUN' if args.dry_run else ('REINDEX' if args.reindex else 'INCREMENTAL')}")
    print("━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n")

    files = discover_files(args.source)

    if not files:
        print("ℹ️   No supported documents found.")
        return

    print(f"🔍  Found {len(files)} document(s).\n")

    if args.dry_run:
        for f in files:
            cat = infer_category(f, args.category)
            print(f"  [{cat:12s}]  {f.relative_to(REPO_ROOT)}")
        return

    # Load existing index
    index = load_index()
    if args.reindex:
        print("♻️   Reindex requested — clearing existing index.\n")
        index["documents"] = {}

    vs = get_vector_store()
    if vs and args.reindex:
        try:
            vs.clear_collection()
            print("   ✅ Vector collection cleared.\n")
        except Exception as exc:
            print(f"   ⚠️  Could not clear collection: {exc}\n")

    # ── Process each file ─────────────────────────────────
    total_chunks  = 0
    skipped       = 0
    failed        = 0

    for file_path in files:
        rel_path = str(file_path.relative_to(REPO_ROOT))
        fhash    = file_hash(file_path)

        # Skip if already indexed and unchanged
        if not args.reindex and rel_path in index["documents"]:
            if index["documents"][rel_path].get("hash") == fhash:
                skipped += 1
                if args.verbose:
                    print(f"   ⏭  Skip (unchanged): {rel_path}")
                continue

        print(f"   📄  Processing: {rel_path}")

        try:
            text = extract_text(file_path)
        except Exception as exc:
            print(f"   ❌  Extract failed: {exc}")
            failed += 1
            continue

        if not text.strip():
            print(f"   ⚠️  Empty or unreadable: {rel_path}")
            failed += 1
            continue

        category = infer_category(file_path, args.category)
        chunks   = chunk_text(text, chunk_size=args.chunk_size)
        print(f"       → {len(chunks)} chunk(s)  |  category: {category}")

        if args.verbose:
            for i, c in enumerate(chunks[:3]):
                print(f"         chunk[{i}]: {c[:80]}…")

        # Add to vector store
        if vs:
            doc_ids = [f"{rel_path}::chunk_{i}" for i in range(len(chunks))]
            metadatas = [
                {
                    "source":    rel_path,
                    "category":  category,
                    "chunk_idx": i,
                    "filename":  file_path.name,
                    "ingested":  datetime.utcnow().isoformat() + "Z",
                }
                for i in range(len(chunks))
            ]
            try:
                vs.add_documents(
                    texts=chunks,
                    ids=doc_ids,
                    metadatas=metadatas,
                )
            except Exception as exc:
                print(f"   ⚠️   Vector store add failed: {exc}")

        # Update index.json
        index["documents"][rel_path] = {
            "hash":        fhash,
            "category":    category,
            "chunks":      len(chunks),
            "char_count":  len(text),
            "ingested_at": datetime.utcnow().isoformat() + "Z",
        }

        total_chunks += len(chunks)

        # Throttle to avoid hammering embedding API
        time.sleep(0.1)

    # ── Save updated index ────────────────────────────────
    save_index(index)

    # ── Summary ───────────────────────────────────────────
    indexed = len(files) - skipped - failed
    print("\n━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━")
    print(f"  ✅  Indexed  : {indexed} file(s)  ({total_chunks} chunks)")
    print(f"  ⏭   Skipped  : {skipped} unchanged")
    print(f"  ❌  Failed   : {failed}")
    print(f"  📁  Index    : {INDEX_FILE}")
    print("━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n")
    print("🧠  Yaazhi now knows your documents. Try asking about them!\n")


# ─────────────────────────────────────────────────────────
# Entry point
# ─────────────────────────────────────────────────────────
if __name__ == "__main__":
    ingest(parse_args())
