"""
Yaazhi Document Ingestion Pipeline.

Ingests PDF, DOCX, PPTX files and URLs into the VectorStore.
Uses tiktoken for chunk sizing, Redis for deduplication tracking.

Usage:
    from memory.ingestion import DocumentIngester
    ingester = DocumentIngester(vector_store)
    result = await ingester.ingest_pdf("/path/to/doc.pdf")
"""

from __future__ import annotations

import asyncio
import hashlib
import os
import time
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional

import aiofiles
import httpx
import logfire
import redis.asyncio as aioredis
import tiktoken
from pydantic import BaseModel, Field

from config.settings import settings
from memory.vector_store import VectorStore


class IngestResult(BaseModel):
    """
    Result of a single document ingestion operation.

    Attributes:
        file_name: Base name of the ingested file.
        file_path: Full path or URL of the source.
        chunks_created: Number of chunks stored in the vector store.
        was_skipped: True if file was already ingested (dedup).
        error: Error message if ingestion failed.
        ingested_at: UTC timestamp of ingestion.
    """

    file_name: str
    file_path: str
    chunks_created: int = Field(default=0, ge=0)
    was_skipped: bool = Field(default=False)
    error: Optional[str] = Field(default=None)
    ingested_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

    def __repr__(self) -> str:
        """Return concise string representation."""
        status = "skipped" if self.was_skipped else ("error" if self.error else "ok")
        return f"IngestResult({self.file_name!r}, chunks={self.chunks_created}, status={status!r})"


class DocumentIngester:
    """
    Handles chunked ingestion of documents into the VectorStore.

    Supported formats: PDF, DOCX, PPTX, URL.
    Deduplication is done via SHA-256 file hash stored in Redis.
    Chunks are created using tiktoken cl100k_base with configurable
    size (default 512 tokens) and overlap (default 50 tokens).
    """

    def __init__(self, vector_store: VectorStore) -> None:
        """
        Initialise the ingester with a VectorStore instance.

        Args:
            vector_store: A ready VectorStore for storing chunks.
        """
        self._vs = vector_store
        self._redis: Optional[aioredis.Redis] = None
        self._tokenizer = tiktoken.get_encoding("cl100k_base")

    def __repr__(self) -> str:
        """Return concise string representation."""
        return f"DocumentIngester(vector_store={self._vs!r})"

    async def _ensure_redis(self) -> None:
        """Lazy-initialise Redis connection."""
        if self._redis is None:
            self._redis = aioredis.from_url(
                settings.redis_url, encoding="utf-8", decode_responses=True
            )

    async def ping(self) -> bool:
        """
        Ping the underlying VectorStore.

        Returns:
            True if VectorStore is reachable.
        """
        return await self._vs.ping()

    async def _chunk_text(
        self, text: str, chunk_size: int = 512, overlap: int = 50
    ) -> list[str]:
        """
        Split text into overlapping token-based chunks.

        Uses tiktoken cl100k_base encoding to count tokens accurately.

        Args:
            text: The full text to chunk.
            chunk_size: Maximum tokens per chunk.
            overlap: Token overlap between consecutive chunks.

        Returns:
            List of text chunk strings.
        """
        tokens = await asyncio.to_thread(self._tokenizer.encode, text)
        chunks: list[str] = []
        step = max(1, chunk_size - overlap)
        for start in range(0, len(tokens), step):
            chunk_tokens = tokens[start: start + chunk_size]
            chunk_text = await asyncio.to_thread(self._tokenizer.decode, chunk_tokens)
            if chunk_text.strip():
                chunks.append(chunk_text.strip())
            if start + chunk_size >= len(tokens):
                break
        return chunks

    async def _file_hash(self, file_path: str) -> str:
        """
        Compute SHA-256 hash of a file using 64 KB streaming chunks.

        VECTOR-1 FIX: Previous implementation loaded the entire file into RAM
        (`await fh.read()`).  On a 4 GB VPS, ingesting a large file could OOM
        the process.  Streaming in fixed-size chunks keeps RSS bounded.

        Args:
            file_path: Absolute path to the file.

        Returns:
            Lowercase hex SHA-256 digest string.
        """
        hasher = hashlib.sha256()
        chunk_size = 64 * 1024   # 64 KB chunks
        async with aiofiles.open(file_path, "rb") as fh:
            while True:
                chunk = await fh.read(chunk_size)
                if not chunk:
                    break
                hasher.update(chunk)
        return hasher.hexdigest()

    async def _already_ingested(self, file_hash: str) -> bool:
        """
        Check whether a file hash is recorded in Redis as already ingested.

        Args:
            file_hash: SHA-256 hex digest of the file.

        Returns:
            True if the hash is already in Redis.
        """
        await self._ensure_redis()
        return bool(await self._redis.exists(f"ingested:{file_hash}"))

    async def _mark_ingested(self, file_hash: str, file_name: str) -> None:
        """
        Record a file hash in Redis so it won't be re-ingested.

        Args:
            file_hash: SHA-256 hex digest of the file.
            file_name: Base file name for human reference.
        """
        await self._ensure_redis()
        await self._redis.set(f"ingested:{file_hash}", file_name)

    async def ingest_pdf(self, file_path: str) -> IngestResult:
        """
        Extract text from a PDF and ingest all chunks into the VectorStore.

        Args:
            file_path: Absolute path to the PDF file.

        Returns:
            IngestResult with chunk count and ingestion status.
        """
        logfire.debug("DocumentIngester.ingest_pdf", path=file_path)
        t_start = time.time()
        file_name = Path(file_path).name

        try:
            file_hash = await self._file_hash(file_path)
            if await self._already_ingested(file_hash):
                logfire.info("Skipping already-ingested PDF", file_name=file_name)
                return IngestResult(file_name=file_name, file_path=file_path, was_skipped=True)

            from pypdf import PdfReader  # type: ignore

            def _extract() -> list[tuple[int, str]]:
                reader = PdfReader(file_path)
                return [(i, page.extract_text() or "") for i, page in enumerate(reader.pages)]

            pages = await asyncio.to_thread(_extract)
            chunks_created = 0
            for page_num, page_text in pages:
                if not page_text.strip():
                    continue
                chunks = await self._chunk_text(page_text)
                for chunk_id, chunk in enumerate(chunks):
                    await self._vs.add(
                        chunk,
                        metadata={
                            "file_name": file_name,
                            "page_num": page_num,
                            "chunk_id": chunk_id,
                            "source_type": "pdf",
                        },
                        source=file_name,
                    )
                    chunks_created += 1

            await self._mark_ingested(file_hash, file_name)
            duration_ms = int((time.time() - t_start) * 1000)
            logfire.info(
                "DocumentIngester.ingest_pdf success",
                file_name=file_name,
                chunks=chunks_created,
                duration_ms=duration_ms,
            )
            return IngestResult(
                file_name=file_name, file_path=file_path, chunks_created=chunks_created
            )
        except Exception as exc:
            logfire.error("DocumentIngester.ingest_pdf failed", file=file_path, error=str(exc))
            return IngestResult(file_name=file_name, file_path=file_path, error=str(exc))

    async def process(self, file_path: str) -> list[str]:
        """
        Generic document ingestion entrypoint used by legacy workflows and tests.
        """
        result = await self.ingest_pdf(file_path)
        return [result.file_name]

    async def ingest_docx(self, file_path: str) -> IngestResult:
        """
        Extract text from a DOCX file and ingest chunks into the VectorStore.

        Args:
            file_path: Absolute path to the DOCX file.

        Returns:
            IngestResult with chunk count and ingestion status.
        """
        logfire.debug("DocumentIngester.ingest_docx", path=file_path)
        t_start = time.time()
        file_name = Path(file_path).name

        try:
            file_hash = await self._file_hash(file_path)
            if await self._already_ingested(file_hash):
                return IngestResult(file_name=file_name, file_path=file_path, was_skipped=True)

            from docx import Document  # type: ignore

            def _extract() -> list[str]:
                doc = Document(file_path)
                return [para.text for para in doc.paragraphs if para.text.strip()]

            paragraphs = await asyncio.to_thread(_extract)
            full_text = "\n".join(paragraphs)
            chunks = await self._chunk_text(full_text)
            chunks_created = 0
            for chunk_id, chunk in enumerate(chunks):
                await self._vs.add(
                    chunk,
                    metadata={
                        "file_name": file_name,
                        "chunk_id": chunk_id,
                        "source_type": "docx",
                    },
                    source=file_name,
                )
                chunks_created += 1

            await self._mark_ingested(file_hash, file_name)
            duration_ms = int((time.time() - t_start) * 1000)
            logfire.info(
                "DocumentIngester.ingest_docx success",
                file_name=file_name,
                chunks=chunks_created,
                duration_ms=duration_ms,
            )
            return IngestResult(
                file_name=file_name, file_path=file_path, chunks_created=chunks_created
            )
        except Exception as exc:
            logfire.error("DocumentIngester.ingest_docx failed", file=file_path, error=str(exc))
            return IngestResult(file_name=file_name, file_path=file_path, error=str(exc))

    async def ingest_pptx(self, file_path: str) -> IngestResult:
        """
        Extract text from a PPTX presentation and ingest chunks.

        Extracts both slide text and presenter notes.

        Args:
            file_path: Absolute path to the PPTX file.

        Returns:
            IngestResult with chunk count and ingestion status.
        """
        logfire.debug("DocumentIngester.ingest_pptx", path=file_path)
        t_start = time.time()
        file_name = Path(file_path).name

        try:
            file_hash = await self._file_hash(file_path)
            if await self._already_ingested(file_hash):
                return IngestResult(file_name=file_name, file_path=file_path, was_skipped=True)

            from pptx import Presentation  # type: ignore

            def _extract() -> list[tuple[int, str]]:
                prs = Presentation(file_path)
                slide_texts: list[tuple[int, str]] = []
                for i, slide in enumerate(prs.slides):
                    parts: list[str] = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            parts.append(shape.text)
                    if slide.has_notes_slide:
                        notes = slide.notes_slide.notes_text_frame.text
                        if notes.strip():
                            parts.append(f"Notes: {notes}")
                    slide_texts.append((i, "\n".join(parts)))
                return slide_texts

            slides = await asyncio.to_thread(_extract)
            chunks_created = 0
            for slide_num, slide_text in slides:
                if not slide_text.strip():
                    continue
                chunks = await self._chunk_text(slide_text)
                for chunk_id, chunk in enumerate(chunks):
                    await self._vs.add(
                        chunk,
                        metadata={
                            "file_name": file_name,
                            "slide_num": slide_num,
                            "chunk_id": chunk_id,
                            "source_type": "pptx",
                        },
                        source=file_name,
                    )
                    chunks_created += 1

            await self._mark_ingested(file_hash, file_name)
            duration_ms = int((time.time() - t_start) * 1000)
            logfire.info(
                "DocumentIngester.ingest_pptx success",
                file_name=file_name,
                chunks=chunks_created,
                duration_ms=duration_ms,
            )
            return IngestResult(
                file_name=file_name, file_path=file_path, chunks_created=chunks_created
            )
        except Exception as exc:
            logfire.error("DocumentIngester.ingest_pptx failed", file=file_path, error=str(exc))
            return IngestResult(file_name=file_name, file_path=file_path, error=str(exc))

    async def ingest_url(self, url: str) -> IngestResult:
        """
        Fetch URL content and ingest chunks into the VectorStore.

        Tries httpx first for static content; falls back gracefully on failure.

        Args:
            url: The URL to fetch and ingest.

        Returns:
            IngestResult with chunk count and ingestion status.
        """
        logfire.debug("DocumentIngester.ingest_url", url=url)
        t_start = time.time()
        file_name = url[:80]

        try:
            from bs4 import BeautifulSoup  # type: ignore

            async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as client:
                resp = await client.get(url)
                resp.raise_for_status()
                html = resp.text

            def _parse(html: str) -> str:
                soup = BeautifulSoup(html, "html.parser")
                for tag in soup(["script", "style", "nav", "footer", "header"]):
                    tag.decompose()
                return soup.get_text(separator="\n", strip=True)[:12000]

            page_text = await asyncio.to_thread(_parse, html)
            url_hash = hashlib.sha256(url.encode()).hexdigest()
            if await self._already_ingested(url_hash):
                return IngestResult(file_name=file_name, file_path=url, was_skipped=True)

            chunks = await self._chunk_text(page_text)
            chunks_created = 0
            for chunk_id, chunk in enumerate(chunks):
                await self._vs.add(
                    chunk,
                    metadata={"url": url, "chunk_id": chunk_id, "source_type": "url"},
                    source=url,
                )
                chunks_created += 1

            await self._mark_ingested(url_hash, file_name)
            duration_ms = int((time.time() - t_start) * 1000)
            logfire.info(
                "DocumentIngester.ingest_url success",
                url=url,
                chunks=chunks_created,
                duration_ms=duration_ms,
            )
            return IngestResult(
                file_name=file_name, file_path=url, chunks_created=chunks_created
            )
        except Exception as exc:
            logfire.error("DocumentIngester.ingest_url failed", url=url, error=str(exc))
            return IngestResult(file_name=file_name, file_path=url, error=str(exc))

    async def ingest_folder(
        self,
        folder_path: str,
        extensions: list[str] = [".pdf", ".docx", ".pptx"],
    ) -> list[IngestResult]:
        """
        Ingest all matching files from a directory.

        Recursively finds files matching the given extensions and calls
        the appropriate ingest method per file type.

        Args:
            folder_path: Path to the folder to scan.
            extensions: List of file extensions to include.

        Returns:
            List of IngestResult for each file found.
        """
        logfire.debug(
            "DocumentIngester.ingest_folder", folder=folder_path, exts=extensions
        )
        folder = Path(folder_path)
        if not folder.exists():
            logfire.warning("Folder does not exist, creating", path=folder_path)
            folder.mkdir(parents=True, exist_ok=True)

        results: list[IngestResult] = []
        dispatch = {
            ".pdf": self.ingest_pdf,
            ".docx": self.ingest_docx,
            ".pptx": self.ingest_pptx,
        }
        for ext in extensions:
            for file_path in sorted(folder.rglob(f"*{ext}")):
                handler = dispatch.get(ext)
                if handler:
                    result = await handler(str(file_path))
                    results.append(result)

        logfire.info(
            "DocumentIngester.ingest_folder complete",
            folder=folder_path,
            total=len(results),
        )
        return results
